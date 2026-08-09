import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_COLOR = RGBColor(15, 23, 42)       # #0F172A Deep Slate Navy
    CARD_BG = RGBColor(30, 41, 59)        # #1E293B Card Slate
    CARD_BORDER = RGBColor(51, 65, 85)    # #334155 Slate Border
    TEXT_BRIGHT = RGBColor(248, 250, 252) # #F8FAFC White Text
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Muted Text
    ACCENT_CYAN = RGBColor(56, 189, 248)  # #38BDF8 Sky/Cyan
    ACCENT_EMERALD = RGBColor(16, 185, 129)# #10B981 Emerald Green
    ACCENT_AMBER = RGBColor(245, 158, 11) # #F59E0B Amber/Gold
    ACCENT_INDIGO = RGBColor(129, 140, 248)# #818CF8 Indigo Accent
    ACCENT_RED = RGBColor(244, 63, 94)    # #F43F5E Rose/Red

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.color.rgb = BG_COLOR
        return bg

    def add_header(slide, slide_num_str, title_text, subtitle_text):
        # Header Badge / Number
        if slide_num_str:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(0.6), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = ACCENT_CYAN
            badge.line.color.rgb = ACCENT_CYAN
            tf = badge.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = slide_num_str
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = BG_COLOR
            p.font.name = "Segoe UI"

        # Title
        title_left = Inches(1.55) if slide_num_str else Inches(0.8)
        tx_box = slide.shapes.add_textbox(title_left, Inches(0.35), Inches(11.0), Inches(0.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_BRIGHT
        p.font.name = "Segoe UI"

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(title_left, Inches(0.85), Inches(11.0), Inches(0.4))
            stf = sub_box.text_frame
            stf.word_wrap = True
            p2 = stf.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_MUTED
            p2.font.name = "Segoe UI"

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    # Accent Glow / Decorative Card
    dec_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    dec_box.fill.solid()
    dec_box.fill.fore_color.rgb = CARD_BG
    dec_box.line.color.rgb = CARD_BORDER

    # Top Tag
    tag_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.8), Inches(4.2), Inches(0.4))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(30, 58, 138)
    tag_box.line.color.rgb = ACCENT_CYAN
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MENTOR & ÜRÜN DEĞERLENDİRME SUNUMU"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"

    # Main Title
    t_box = s1.shapes.add_textbox(Inches(2.0), Inches(2.4), Inches(9.3), Inches(1.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BIST BOT"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_BRIGHT
    p.font.name = "Segoe UI"

    # Subtitle
    sub_box = s1.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9.3), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Borsa İstanbul Verisini Karara Dönüştüren Açıklanabilir Sinyal Motoru"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"

    # Description
    desc_box = s1.shapes.add_textbox(Inches(2.0), Inches(4.6), Inches(9.3), Inches(0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Çoklu zaman dilimi taraması, şeffaf skorlama, sert veto kuralları, otonom risk yönetimi ve walk-forward doğrulama altyapısı."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Segoe UI"

    # Footer Info
    foot_box = s1.shapes.add_textbox(Inches(2.0), Inches(5.5), Inches(9.3), Inches(0.5))
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hazırlayan: Mustafa Ali Ertuğrul  |  Ağustos 2026"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 2: PROBLEM & TARGET AUDIENCE
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "02", "PROBLEM VE HEDEF KULLANICI KİTLESİ", "Sorun veri eksikliği değil; karar disiplini, açıklanabilirlik ve parçalı araç problemidir.")

    # Left Card: Bireysel Yatırımcı
    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(4.3))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = CARD_BORDER

    t1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(0.4))
    p = t1.text_frame.paragraphs[0]
    p.text = "👤 Bireysel Yatırımcı"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER
    p.font.name = "Segoe UI"

    t1_body = s2.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(3.4))
    tf = t1_body.text_frame
    tf.word_wrap = True
    items1 = [
        ("Sinyal Kirliliği:", " Çok sayıda indikatörü eş zamanlı takip ederken çelişkili sinyaller arasında kararsız kalır."),
        ("Kara Kutu Şüpheciliği:", " Üretilen AL/SAT alarmlarının arkasındaki mantığı ve risk nedenini göremez."),
        ("Disiplinsiz İşlem:", " Stop-loss, hedef fiyat ve pozisyon büyüklüğü planlamadan duygusal işlemler yapar.")
    ]
    for bold_prefix, text in items1:
        p = tf.add_paragraph()
        p.font.size = Pt(13)
        p.font.name = "Segoe UI"
        run1 = p.add_run()
        run1.text = "• " + bold_prefix
        run1.font.bold = True
        run1.font.color.rgb = TEXT_BRIGHT
        run2 = p.add_run()
        run2.text = text
        run2.font.color.rgb = TEXT_MUTED

    # Right Card: Operatör / Analist
    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.7), Inches(4.3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = CARD_BORDER

    t2 = s2.shapes.add_textbox(Inches(7.033), Inches(1.6), Inches(5.3), Inches(0.4))
    p = t2.text_frame.paragraphs[0]
    p.text = "📊 Disiplinli Analist & Operatör"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"

    t2_body = s2.shapes.add_textbox(Inches(7.033), Inches(2.1), Inches(5.3), Inches(3.4))
    tf = t2_body.text_frame
    tf.word_wrap = True
    items2 = [
        ("Parçalı Araç Yükü:", " Tarama, grafik inceleme, risk hesabı ve bildirim için ayrı platformlar kullanmak zorundadır."),
        ("Tekilleştirilemeyen Testler:", " Strateji fikirlerini geriye dönük (backtest) ve canlı öncesi (paper trade) tutarlı doğrulamakta zorlanır."),
        ("Ölçeklenme Engeli:", " BIST evrenindeki tüm hisseleri üst zaman dilimi (MTF) uyumuyla manuel tarayamaz.")
    ]
    for bold_prefix, text in items2:
        p = tf.add_paragraph()
        p.font.size = Pt(13)
        p.font.name = "Segoe UI"
        run1 = p.add_run()
        run1.text = "• " + bold_prefix
        run1.font.bold = True
        run1.font.color.rgb = TEXT_BRIGHT
        run2 = p.add_run()
        run2.text = text
        run2.font.color.rgb = TEXT_MUTED

    # Bottom Value Proposition Banner
    b_banner = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.9), Inches(11.733), Inches(1.1))
    b_banner.fill.solid()
    b_banner.fill.fore_color.rgb = RGBColor(20, 83, 45) # Dark Emerald Green
    b_banner.line.color.rgb = ACCENT_EMERALD

    tb = b_banner.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = "💡 BIST BOT’UN İDDİASI VE DEĞER ÖNERMESİ"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.font.name = "Segoe UI"

    p2 = tb.add_paragraph()
    p2.text = "Tarama  ➔  Açıklanabilir Skor  ➔  Risk Planı  ➔  Telegram Bildirimi  ➔  Paper Trade  ➔  Walk-Forward Backtest akışını uçtan uca TEK otonom sistemde birleştirmektir."
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_BRIGHT
    p2.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 3: END-TO-END PIPELINE
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "03", "OTONOM KARAR VE YÜRÜTME HATTI (PIPELINE)", "Her BIST hissesi 5 aşamalı standart, disiplinli ve izlenebilir bir süreçten geçer.")

    steps = [
        ("1. VERİ AKIŞI", "BIST Evrensel Verisi", ["• BIST 100 / Tüm Hisseler", "• OHLCV Fiyat & Hacim", "• Çoklu Zaman Dilimi (MTF)", "  (Günlük, Saatlik, 15dk)"], ACCENT_CYAN),
        ("2. İNDİKATÖRLER", "Teknik Metrikler", ["• RSI, MACD, Bollinger", "• SMA/EMA Kesişimleri", "• ADX Trend Gücü", "• OBV & Fiyat-Hacim"], ACCENT_INDIGO),
        ("3. SKOR & VETO", "Karar Motoru", ["• 4 Kanıt Grubu Puanı", "• -100 ile +100 Skor", "• Sert Veto Filtreleri", "• Şeffaf Açıklama Listesi"], ACCENT_AMBER),
        ("4. RİSK YÖNETİMİ", "Pozisyon Planlaması", ["• Dinamik Stop-Loss", "• Kar Al (Take-Profit)", "• ATR Volatilitesi", "• Risk/Ödül (R/R) Kontrolü"], ACCENT_EMERALD),
        ("5. İZLEME", "Bildirim & Doğrulama", ["• Telegram Anlık Bildirim", "• Otonom Paper Trade", "• Canlı Takip & Loglama", "• Walk-Forward Kaydı"], ACCENT_CYAN)
    ]

    card_w = Inches(2.2)
    card_h = Inches(5.0)
    start_x = Inches(0.8)
    gap_x = Inches(0.18)

    for i, (num_title, sub_title, bullet_list, accent_clr) in enumerate(steps):
        cx = start_x + i * (card_w + gap_x)
        cy = Inches(1.6)

        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Header Box inside card
        hdr_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.1), cy + Inches(0.1), card_w - Inches(0.2), Inches(0.85))
        hdr_box.fill.solid()
        hdr_box.fill.fore_color.rgb = BG_COLOR
        hdr_box.line.color.rgb = accent_clr

        tf = hdr_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num_title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = accent_clr
        p.font.name = "Segoe UI"

        p_sub = tf.add_paragraph()
        p_sub.text = sub_title
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = TEXT_BRIGHT
        p_sub.font.name = "Segoe UI"

        # Content Box
        ct_box = s3.shapes.add_textbox(cx + Inches(0.1), cy + Inches(1.05), card_w - Inches(0.2), Inches(3.7))
        tf = ct_box.text_frame
        tf.word_wrap = True
        for b_idx, bullet_line in enumerate(bullet_list):
            if b_idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet_line
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_BRIGHT if bullet_line.startswith("•") else TEXT_MUTED
            p.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 4: 4 EVIDENCE GROUPS
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "04", "KARAR MEKANİZMASI: 4 TEMEL KANIT GRUBU", "Sinyaller tek bir göstergeye değil, 4 bağımsız teknik boyuta dayanır.")

    ev_groups = [
        ("⚡ MOMENTUM KANITLARI", "Aşırı Alım / Satım & Hız", [
            "RSI (Göreceli Güç Endeksi)",
            "Stochastic Osilatör Kesişimleri",
            "CCI (Emtia Kanal Endeksi)",
            "Aşırı Alım/Satım Bölge Teyitleri"
        ], ACCENT_CYAN),
        ("📈 TREND & REJİM KANITLARI", "Yön Gücü & Zaman Dilimi", [
            "SMA / EMA Hareketli Ortalama Kesişimleri",
            "MACD Sinyali & Histogram Trendi",
            "ADX + DI Trend Gücü (Trend vs Yatay)",
            "Üst Zaman Dilimi (MTF) Yön Uyumu"
        ], ACCENT_EMERALD),
        ("📊 HACİM & AKIŞ KANITLARI", "Katılım & Hacim Doğrulaması", [
            "Hacim Katı (Volume Spike Tespiti)",
            "OBV (Denge İşlem Hacmi) Trendi",
            "Fiyat - Hacim Uyumsuzluğu (Divergence)",
            "Nakit Akışı ve Para Giriş Teyidi"
        ], ACCENT_AMBER),
        ("🛡️ YAPI & RİSK KANITLARI", "Volatolite & Yapısal Seviyeler", [
            "Bollinger Bant Konumu & Daralması",
            "Statik Destek / Direnç Seviyeleri",
            "ATR (Ortalama Gerçek Aralık) Volatilitesi",
            "Risk / Ödül (R/R) & Likidite İncelemesi"
        ], ACCENT_INDIGO)
    ]

    card_w = Inches(5.6)
    card_h = Inches(2.5)

    positions = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.933), Inches(1.6)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.933), Inches(4.35))
    ]

    for i, (g_title, g_sub, g_bullets, accent_clr) in enumerate(ev_groups):
        cx, cy = positions[i]
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Header Box
        hdr_box = s4.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.15), card_w - Inches(0.4), Inches(0.5))
        tf = hdr_box.text_frame
        p = tf.paragraphs[0]
        p.text = g_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_clr
        p.font.name = "Segoe UI"

        # Bullets Box
        b_box = s4.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.65), card_w - Inches(0.4), Inches(1.7))
        tf = b_box.text_frame
        tf.word_wrap = True
        for bullet in g_bullets:
            p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_BRIGHT
            p.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 5: EXPLAINABLE SCORING & VETO FILTERS
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "05", "ŞEFFAF SKORLAMA VE SERT VETO FİLTRELERİ", "Sinyaller kara kutu değildir; birleşik skor ve katı veto kuralları ile üretilir.")

    # Top Card: Score Spectrum
    top_c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.0))
    top_c.fill.solid()
    top_c.fill.fore_color.rgb = CARD_BG
    top_c.line.color.rgb = CARD_BORDER

    tf = top_c.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 BİRLEŞİK SKOR SPEKTRUMU (-100  ←  0  →  +100)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"

    # Score breakdown boxes
    scores_info = [
        ("🚨 GÜÇLÜ SAT", "Skor ≤ -48", ACCENT_RED),
        ("🔴 SAT", "Skor ≤ -20", RGBColor(225, 29, 72)),
        ("🟡 ZAYIF SAT", "Skor ≤ -8", ACCENT_AMBER),
        ("⚪ BEKLE", "-8 < Skor < 8", TEXT_MUTED),
        ("🟢 ZAYIF AL", "Skor ≥ 8", RGBColor(134, 239, 172)),
        ("🟢 AL", "Skor ≥ 20", RGBColor(34, 197, 94)),
        ("💰 GÜÇLÜ AL", "Skor ≥ 48", ACCENT_EMERALD)
    ]

    sb_w = Inches(1.55)
    for j, (s_name, s_val, s_clr) in enumerate(scores_info):
        sx = Inches(1.0) + j * (sb_w + Inches(0.1))
        sy = Inches(2.2)
        sb = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sb_w, Inches(1.0))
        sb.fill.solid()
        sb.fill.fore_color.rgb = BG_COLOR
        sb.line.color.rgb = s_clr
        stf = sb.text_frame
        stf.word_wrap = True
        p1 = stf.paragraphs[0]
        p1.text = s_name
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = s_clr
        p1.font.name = "Segoe UI"
        p2 = stf.add_paragraph()
        p2.text = s_val
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BRIGHT
        p2.font.name = "Segoe UI"

    # 3 Bottom Veto Cards
    vetos = [
        ("⛔ 1. Yön ve Rejim Çelişkisi", "Trend ile Tetikleyici Uyumsuzluğu", "Ana trend yönü (örneğin düşen trend) ile anlık indikatör sinyali çelişiyorsa sinyal skoru otomatik olarak bastırılır veya veto edilir.", ACCENT_AMBER),
        ("⛔ 2. Aşırı Uzama (Chase Risk)", "Direnç Yakınlığı & Doygunluk", "Fiyat üst Bollinger bandına veya kritik dirence çok yakınsa chase (fiyatı kovalama) riski tespiti yapılır; skora ceza verilir.", ACCENT_RED),
        ("⛔ 3. Yapısal Zayıflık & Likidite", "ATR Uç Noktaları & Hacim Eksikliği", "Düşük işlem hacmi veya volatilite patlaması yaşanan ortamlarda sahte kırılmaları önlemek adına sinyal yürütmesi engellenir.", ACCENT_INDIGO)
    ]

    vw = Inches(3.7)
    vh = Inches(3.2)
    for i, (v_title, v_sub, v_desc, v_clr) in enumerate(vetos):
        vx = Inches(0.8) + i * (vw + Inches(0.3))
        vy = Inches(3.8)

        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, vx, vy, vw, vh)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = v_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = v_clr
        p.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = v_sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_BRIGHT
        p2.font.name = "Segoe UI"

        p3 = tf.add_paragraph()
        p3.text = "\n" + v_desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = TEXT_MUTED
        p3.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 6: ML & SELF-TRAINING
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "06", "OTONOM GERİ BESLEME VE MODEL KALİBRASYONU", "Zaman serisi hassasiyetini koruyan kontrollü makine öğrenmesi döngüsü.")

    ml_steps = [
        ("01", "Sinyal Geçmişi & Özellik Kaydı", "Sinyal anındaki Skor, ADX, RSI, Hacim Katı, ATR, R/R Oranı ve Trend ölçekleri veritabanına eksiksiz kaydedilir.", ACCENT_CYAN),
        ("02", "İleri Dönem Etiketleme", "Üretilen sinyaller, N bar sonraki gerçekleşen getiri performansına göre Başarılı / Başarısız olarak etiketlenir.", ACCENT_INDIGO),
        ("03", "TimeSeriesSplit Eğitimi", "Gelecek bilgisinin sızmasını (Look-ahead bias) önlemek için TimeSeriesSplit ile XGBoost / Lojistik Regresyon eğitilir.", ACCENT_AMBER),
        ("04", "Canlı Olasılık Kalibrasyonu", "Ham model çıktısı Platt Scaling / Isotonic yöntemle kalibre edilir; üretimde güven olasılığı skoru üretir.", ACCENT_EMERALD)
    ]

    card_w = Inches(2.75)
    card_h = Inches(4.8)
    for i, (num_str, title_str, desc_str, clr_val) in enumerate(ml_steps):
        cx = Inches(0.8) + i * (card_w + Inches(0.24))
        cy = Inches(1.7)

        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Large Number Badge
        nb = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.2), cy + Inches(0.2), Inches(0.8), Inches(0.6))
        nb.fill.solid()
        nb.fill.fore_color.rgb = BG_COLOR
        nb.line.color.rgb = clr_val
        np = nb.text_frame.paragraphs[0]
        np.text = num_str
        np.alignment = PP_ALIGN.CENTER
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = clr_val
        np.font.name = "Segoe UI"

        tf = s6.shapes.add_textbox(cx + Inches(0.15), cy + Inches(1.0), card_w - Inches(0.3), Inches(3.6)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_str
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_BRIGHT
        p.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = "\n" + desc_str
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 7: BACKTEST & TEST RESULTS
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "07", "WALK-FORWARD VE MALİYET STRESİ ALTINDA TEST SONUÇLARI", "Geriye dönük doğrulama sonuçları umut vericidir; risk sınırları şeffafça tanımlanmıştır.")

    # Left Highlight Box 1
    hb1 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.3))
    hb1.fill.solid()
    hb1.fill.fore_color.rgb = CARD_BG
    hb1.line.color.rgb = ACCENT_EMERALD

    tf1 = hb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "8 / 30 HİSSE"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.font.name = "Segoe UI"

    p2 = tf1.add_paragraph()
    p2.text = "Komisyon, Kayma (Slippage) ve Maliyet Stresi Altında 'Robust' (Dayanıklı) Kalan Hisse Sayısı"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_BRIGHT
    p2.font.name = "Segoe UI"

    # Right Highlight Box 2
    hb2 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.6), Inches(5.6), Inches(2.3))
    hb2.fill.solid()
    hb2.fill.fore_color.rgb = CARD_BG
    hb2.line.color.rgb = ACCENT_CYAN

    tf2 = hb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "%26,7"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"

    p2 = tf2.add_paragraph()
    p2.text = "İlk Walk-Forward Test Evrenindeki Dayanıklı (Robust) Strateji Başarı Payı"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_BRIGHT
    p2.font.name = "Segoe UI"

    # Bottom Risk & Boundaries Card
    rb = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(11.733), Inches(2.7))
    rb.fill.solid()
    rb.fill.fore_color.rgb = CARD_BG
    rb.line.color.rgb = ACCENT_AMBER

    tf = rb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ ŞEFFAF RİSK VE SINIRLILIK İTİRAFI (PRODUCTION RISKS)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER
    p.font.name = "Segoe UI"

    r_items = [
        ("Tarihsel Veri Kısıtı:", " Elde edilen sonuçlar geçmiş piyasa rejimlerine ve kısıtlı işlem sayılarına dayanmaktadır; gelecekteki performansı doğrudan garanti etmez."),
        ("15 Dakika Veri Gecikmesi:", " Üretim ortamında 15 dk gecikmeli veri kullanımı ani volatilitede kayma (slippage) riskini artırır."),
        ("Survivorship Bias:", " Test evrenindeki hisseler güncel listeden seçildiği için geçmişte borsa kotundan çıkan hisselerin süzülmesi canlı ortamda dikkate alınmalıdır.")
    ]
    for bold_prefix, text in r_items:
        p = tf.add_paragraph()
        p.font.size = Pt(12)
        p.font.name = "Segoe UI"
        run1 = p.add_run()
        run1.text = "• " + bold_prefix
        run1.font.bold = True
        run1.font.color.rgb = TEXT_BRIGHT
        run2 = p.add_run()
        run2.text = text
        run2.font.color.rgb = TEXT_MUTED


    # ==========================================
    # SLIDE 8: DIFFERENTIATION MATRIX
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "08", "MEVCUT ÇÖZÜMLERDEN FARKIMIZ", "Farkımız 'daha çok indikatör' değil; uçtan uca izlenebilir ve doğrulanabilir karar sistemidir.")

    # Table creation
    rows, cols = 5, 3
    left, top, width, height = Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0)
    table_shape = s8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(4.3)
    table.columns[2].width = Inches(4.633)

    headers = ["KARAR BOYUTU", "GELENEKSEL ARAÇLAR / ALARMLAR", "BIST BOT SİSTEM MİMARİSİ"]
    for j, h_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Segoe UI"

    table_data = [
        ("Sinyal Üretimi", "Tek hisse / Tek grafik seviyesinde basit indikatör alarmı", "Tüm BIST evreninde çoklu zaman dilimi (MTF) & piyasa rejimi taraması"),
        ("Açıklanabilirlik", "Sadece 'AL' ya da 'SAT' üreten şeffaf olmayan kara kutu", "Her sinyalin alt bileşen puanları, 4 kanıt grubu ve sert veto gerekçeleri"),
        ("Risk & Yürütme", "Risk hesabı tamamen kullanıcının anlık insiyatifine bırakılır", "Dinamik Stop / Hedef seviyeleri, ATR pozisyon boyutu ve risk limitleri"),
        ("Doğrulama Döngüsü", "Tekilleştirilmiş, maliyetsiz ve basit geriye dönük testler", "Walk-forward doğrulama, maliyet/kayma stresi ve otonom paper trade")
    ]

    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(12)
            p.font.name = "Segoe UI"
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_AMBER
            elif j == 1:
                p.font.color.rgb = TEXT_MUTED
            else:
                p.font.bold = True
                p.font.color.rgb = TEXT_BRIGHT


    # ==========================================
    # SLIDE 9: PRODUCT VISION & ROADMAP
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "09", "ÜRÜN VİZYONU VE YOL HARİTASI", "Araştırma prototipinden kurumsal ve güvenilir karar altyapısına geçiş planı.")

    phases = [
        ("BUGÜN", "Mevcut Durum / Prototip", [
            "Açıklanabilir BIST tarama motoru",
            "Telegram anlık bildirim sistemi",
            "Otonom paper trade altyapısı",
            "Walk-forward backtest modülü",
            "Streamlit & Flask operatör panelleri"
        ], ACCENT_CYAN),
        ("0 – 6 AY", "Yakın Gelecek / Üretim", [
            "BIST canlı resmi veri akışı (Websocket)",
            "Model sürüklenmesi (Drift) takibi & uyarıları",
            "Kapalı beta kullanıcı testi ve geri bildirim",
            "Gerçek zamanlı performans dashboard'u",
            "Cloud Run / Cloud SQL güvenilir dağıtım"
        ], ACCENT_EMERALD),
        ("6 – 18 AY", "Uzun Vadeli Vizyon", [
            "AlgoLab / Broker entegrasyonu (Otomatik Emir)",
            "Kişiselleştirilmiş kullanıcı risk profilleri",
            "Çoklu varlık genişlemesi (Kripto & Global)",
            "Kurumsal API ve strateji marketplace",
            "Mobil uygulama arayüzü"
        ], ACCENT_INDIGO)
    ]

    card_w = Inches(3.7)
    card_h = Inches(5.0)
    for i, (p_title, p_sub, p_bullets, accent_clr) in enumerate(phases):
        cx = Inches(0.8) + i * (card_w + Inches(0.3))
        cy = Inches(1.6)

        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Phase Header Box
        ph_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.15), cy + Inches(0.15), card_w - Inches(0.3), Inches(0.9))
        ph_box.fill.solid()
        ph_box.fill.fore_color.rgb = BG_COLOR
        ph_box.line.color.rgb = accent_clr

        tf = ph_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = p_title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_clr
        p.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = p_sub
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BRIGHT
        p2.font.name = "Segoe UI"

        # Bullets Box
        b_box = s9.shapes.add_textbox(cx + Inches(0.15), cy + Inches(1.15), card_w - Inches(0.3), Inches(3.6))
        tf = b_box.text_frame
        tf.word_wrap = True
        for bullet in p_bullets:
            p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_BRIGHT
            p.font.name = "Segoe UI"


    # ==========================================
    # SLIDE 10: MENTOR FEEDBACK & CONCLUSION
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_header(s10, "10", "MENTOR DEĞERLENDİRMESİ VE ODAK NOKTALARI", "Projenin bir sonraki aşamaya taşınmasında kritik tavsiye ve rehberlik alanları.")

    questions = [
        ("1️⃣ Hedef Kullanıcı & Senaryo", "Bireysel yatırımcı paneli mi yoksa analist odaklı otonom karar destek motoru mu önceliklendirilmeli?", ACCENT_CYAN),
        ("2️⃣ Canlı Pilot Standartları", "Canlı emre geçiş öncesinde gerekli walk-forward kanıt eşiği ve risk sınırları ne olmalı?", ACCENT_AMBER),
        ("3️⃣ Veri & Entegrasyon Yol Haritası", "Lisanslı veri tedariki ve aracı kurum (broker) entegrasyonunda en verimli ve sürdürülebilir strateji hangisidir?", ACCENT_EMERALD)
    ]

    card_w = Inches(3.7)
    card_h = Inches(3.6)
    for i, (q_title, q_desc, q_clr) in enumerate(questions):
        cx = Inches(0.8) + i * (card_w + Inches(0.3))
        cy = Inches(1.6)

        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = q_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = q_clr
        p.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = "\n" + q_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_BRIGHT
        p2.font.name = "Segoe UI"

    # Bottom Mission Statement Box
    b_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.4))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = RGBColor(30, 58, 138)
    b_box.line.color.rgb = ACCENT_CYAN

    tf = b_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 TEMEL MİSYONUMUZ VE AMACIMIZ"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = "BIST BOT’un amacı al-sat tavsiyesi vermek değil; yatırım kararlarını ölçülebilir, açıklanabilir, risk odaklı ve disiplinli bir altyapıya kavuşturmaktır."
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_BRIGHT
    p2.font.name = "Segoe UI"

    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    out = r"C:\Users\Ali\OneDrive\Masaüstü\bist_bot\BIST_BOT_Mentor_Sunumu.pptx"
    build_presentation(out)
